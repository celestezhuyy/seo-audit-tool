import streamlit as st
import time
import pandas as pd
import requests
import hashlib
import re
import urllib3
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from io import BytesIO
import socket

# --- Level 0: 页面基础配置 ---
st.set_page_config(
    page_title="NextGen SEO Auditor",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 禁用 SSL 警告
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# 尝试导入 pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
except ImportError:
    st.error("Missing dependencies! Please add 'python-pptx' to requirements.txt.")
    st.stop()

# --- Level 1: 基础工具函数 ---
def is_valid_url(url):
    try:
        result = urlparse(url)
        return all([result.scheme, result.netloc])
    except:
        return False

def get_content_hash(text):
    return hashlib.md5(text.encode('utf-8')).hexdigest()

def estimate_pixel_width(text, font_size=18):
    if not text: return 0
    width = 0
    for char in text:
        if ord(char) > 127: 
            width += font_size
        elif char.isupper():
            width += font_size * 0.7 
        else:
            width += font_size * 0.55 
    return width

def get_browser_headers():
    return {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
        'Connection': 'keep-alive',
    }

def fetch_psi_data(url, api_key):
    if not api_key: return None
    endpoint = f"https://www.googleapis.com/pagespeedonline/v5/runPagespeed?url={url}&key={api_key}&strategy=mobile"
    try:
        response = requests.get(endpoint, timeout=30)
        if response.status_code == 200:
            data = response.json()
            crux = data.get('loadingExperience', {}).get('metrics', {})
            if not crux: return {"error": "No CrUX data available"}
            return {
                "LCP": crux.get('LARGEST_CONTENTFUL_PAINT_MS', {}).get('percentile', 0) / 1000,
                "CLS": crux.get('CUMULATIVE_LAYOUT_SHIFT_SCORE', {}).get('percentile', 0) / 100,
                "INP": crux.get('INTERACTION_TO_NEXT_PAINT', {}).get('percentile', 0),
                "FCP": crux.get('FIRST_CONTENTFUL_PAINT_MS', {}).get('percentile', 0) / 1000,
            }
        else: return {"error": f"API Error: {response.status_code}"}
    except Exception as e: return {"error": str(e)}

def check_server_location(url):
    try:
        domain = urlparse(url).netloc
        ip = socket.gethostbyname(domain)
        response = requests.get(f"http://ip-api.com/json/{ip}", timeout=3)
        if response.status_code == 200:
            data = response.json()
            return data.get("countryCode", "Unknown"), data.get("country", "Unknown")
    except:
        pass
    return None, None

# --- Level 2: 排序与配置常量 ---
CATEGORY_ORDER = ["access", "indexability", "technical", "content", "image_ux", "cwv_performance"]
SEVERITY_ORDER = {"Critical": 0, "High": 1, "Medium": 2, "Low": 3}

ISSUE_PRIORITY_LIST = [
    "no_robots", "robots_bad_rule", "robots_quality_issue", "baidu_robots_missing", "robots_no_sitemap", "no_sitemap", "sitemap_invalid",
    "http_5xx", "http_4xx", "soft_404", "http_3xx",
    "server_not_in_china",
    "duplicate", "missing_canonical", "hreflang_invalid", "hreflang_no_default", "missing_hreflang",
    "missing_viewport", "missing_jsonld", "js_links", "url_underscore", "url_uppercase",
    "missing_baidu_stats", "missing_baidu_verify", "missing_applicable_device", "missing_no_transform", "missing_icp", "content_not_chinese",
    "missing_title", "short_title", "long_title", "missing_desc", "short_desc", "missing_h1", "missing_keywords", 
    "no_favicon", "missing_alt", "alt_bad_quality", "anchor_bad_quality", 
    "lcp_issue", "inp_issue", "cls_issue", "fcp_issue", "fcp_baidu_issue", "cls_risk"
]

def get_issue_priority(issue_id):
    try: return ISSUE_PRIORITY_LIST.index(issue_id)
    except ValueError: return 999 

# --- Level 3: 国际化字典 ---
TRANSLATIONS = {
    "zh": {
        "sidebar_title": "🔍 AuditAI Pro",
        "sidebar_caption": "旗舰审计版 v15.3",
        "nav_label": "功能导航",
        "nav_options": ["输入网址", "仪表盘", "数据矩阵", "PPT 生成器"],
        "lang_label": "语言 / Language",
        "clear_data": "清除数据并重置",
        "cache_info": "已缓存 {} 个页面",
        "sitemap_status_title": "Sitemap 状态:",
        "sitemap_found_href": "✅ 发现 Hreflang 配置", 
        "sitemap_no_href": "⚠️ 未发现 Hreflang",     
        "sitemap_missing": "❌ 未找到 Sitemap",
        "psi_settings": "Google PSI API 设置 (推荐)",
        "psi_api_key_label": "输入 Google PageSpeed API Key",
        "psi_api_help": "建议填入以获取 LCP/CLS/INP 真实数据。留空则只进行代码审计。",
        "psi_list_url_label": "产品列表页 URL (可选)",
        "psi_detail_url_label": "产品详情页 URL (可选)",
        "psi_get_key": "没有 API Key? [点击这里免费申请](https://developers.google.com/speed/docs/insights/v5/get-started)",
        "psi_fetching": "正在调用 Google API 获取 {} 数据...",
        "psi_success": "成功获取真实用户数据！",
        "psi_error": "API 调用失败或无 CrUX 数据",
        "input_header": "开始深度审计",
        "input_info": "说明: v15.3 修复了缩进错误，包含百度 SEO、手动 Sitemap 解析等增强功能。",
        "input_label": "输入目标网址 (首页)",
        "input_placeholder": "https://example.com",
        "max_pages_label": "最大爬取页面数",
        "adv_settings": "高级设置 (Advanced Settings)", 
        "check_robots_label": "检查并遵循 Robots.txt 规则", 
        "crawl_sitemap_label": "自动抓取 Robots.txt 中的 Sitemap", 
        "baidu_mode_label": "启用百度 SEO 审计模式", 
        "allow_subdomains_label": "允许抓取子域名 (如 blog.site.com)",
        "allow_outside_folder_label": "允许抓取父级目录 (如从 /en/ 开始抓取 /fr/)",
        "manual_sitemaps": "手动 Sitemap 地址 (每行一个, 补充用)", 
        "manual_pages_label": "手动添加页面列表 (每行一个 URL)",
        "sitemap_content_label": "粘贴 Sitemap XML 内容 (直接解析)",
        "start_btn": "开始深度爬取",
        "error_url": "网址格式错误",
        "spinner_crawl": "正在执行深度审计 (Max {} pages)...", 
        "error_no_data": "未能爬取到任何页面。原因: {}", 
        "success_audit": "审计完成！共分析 {} 个页面。",
        "dashboard_header": "执行摘要 (Executive Summary)",
        "warn_no_data": "暂无数据。",
        "kpi_health": "网站健康度",
        "kpi_pages": "已分析页面",
        "kpi_issues": "发现问题总数",
        "kpi_critical": "严重问题",
        "chart_issues": "问题类型分布",
        "chart_no_issues": "未发现明显问题。",
        "chart_status": "HTTP Status Codes",
        "cwv_title": "首页核心 Web 指标 (Core Web Vitals) - 真实数据",
        "cwv_source": "数据来源: Google Chrome User Experience Report (CrUX) - 仅首页",
        "matrix_header": "爬取数据明细 (Big Sheet)",
        "download_csv": "下载 CSV 报告",
        "ppt_header": "演示文稿预览 (Pitch Deck Mode)",
        "ppt_success_no_issues": "无严重问题。",
        "ppt_download_header": "📥 导出报告",
        "ppt_info": "说明：生成的 PPT 已优化为 16:9 宽屏，包含增强版可视化预览。",
        "ppt_btn": "生成并下载美化版 .pptx",
        "ppt_preview_header": "网页版预览",
        "ppt_slide_title": "问题类型:",
        "ppt_category": "分类:",
        "ppt_severity": "严重程度:",
        "ppt_impact": "影响范围:",
        "ppt_impact_desc": "在已爬取样本中发现 **{}** 个页面。",
        "ppt_desc": "🔴 问题描述:",
        "ppt_business_impact": "📉 商业与 SEO 影响:", 
        "ppt_sugg": "✅ 修复建议:",
        "ppt_examples": "🔍 受影响页面示例:",
        "ppt_prev": "⬅️ 上一页",
        "ppt_next": "下一页 ➡️",
        "cat_access": "可访问性与索引 (Access & Indexing)",
        "cat_indexability": "索引规范性 (Indexability)",
        "cat_technical": "技术与架构 (Technical SEO)",
        "cat_content": "页面内容 (On-Page Content)",
        "cat_image_ux": "用户体验与资源 (UX & Assets)",
        "cat_cwv_performance": "核心性能指标 (Core Web Vitals)",
        "ppt_cover_title": "SEO 深度技术审计报告",
        "ppt_cover_sub": "Generated by AuditAI Pro v15.2",
        "ppt_slide_desc_title": "深度分析",
        "ppt_slide_count_title": "样本中受影响页面数: {} 个",
        "ppt_slide_ex_title": "受影响页面示例", 
        "ppt_slide_sugg_title": "💡 修复建议:",
        "serp_sim_title": "Google 搜索结果模拟 (SERP):",
        "rich_sim_title": "富媒体结果模拟 (Rich Results):",
        "code_sim_title": "Code Snippet:",
        "visual_sim_title": "视觉体验模拟:",
        "cwv_sim_title": "CWV 性能仪表盘 (Performance):",
        "lcp_issue": "LCP (最大内容绘制) 超标", "lcp_issue_desc": "LCP 时间为 {:.2f}s (目标 <2.5s)。页面主要内容加载过于缓慢。", "lcp_issue_impact": "LCP 是 Google 核心排名因素。加载缓慢会导致用户跳出率飙升，并直接降低在移动端的搜索排名。", "lcp_issue_sugg": "压缩图片体积（使用 WebP），使用 CDN 分发内容，推迟非关键 JS 执行，并预加载 LCP 关键元素。",
        "cls_issue": "CLS (累积布局偏移) 超标", "cls_issue_desc": "页面加载过程中元素发生意外位移 (Score > 0.1)。", "cls_issue_impact": "作为核心排名因素，布局不稳定会导致用户误触广告或按钮，严重损害品牌信誉和用户体验。", "cls_issue_sugg": "为所有图片和视频元素指定明确的宽度和高度属性，避免在顶部动态插入内容。",
        "inp_issue": "INP (交互到绘制延迟) 超标", "inp_issue_desc": "用户点击按钮后，页面响应延迟超过 200ms。", "inp_issue_impact": "Google 新引入的核心指标。高延迟会让用户觉得网站“卡顿”或无响应，严重影响转化率。", "inp_issue_sugg": "减少主线程阻塞，将长任务 (Long Tasks) 拆分为小任务，并优化复杂的 JavaScript 事件处理逻辑。",
        "fcp_issue": "FCP (首次内容绘制) 缓慢", "fcp_issue_desc": "FCP 时间为 {:.2f}s (目标 <1.8s)。用户看到页面第一个内容的时间过长。", "fcp_issue_impact": "FCP 慢会让用户感觉服务器响应迟钝，直接增加跳出率。", "fcp_issue_sugg": "优化服务器响应时间 (TTFB)，消除阻塞渲染的 CSS/JS 资源。",
        "fcp_baidu_issue": "FCP 不满足百度秒开要求", "fcp_baidu_issue_desc": "FCP 时间为 {:.2f}s (百度目标 < 2.0s)。", "fcp_baidu_issue_impact": "百度对落地页首屏加载速度有严格要求，超过2秒将导致移动搜索降权。", "fcp_baidu_issue_sugg": "使用国内 CDN，精简首屏资源，确保 FCP < 2s。",
        "no_robots": "缺失 Robots.txt", "no_robots_desc": "无法访问根目录的 robots.txt 文件，或者服务器返回错误状态码。", "no_robots_impact": "爬虫可能抓取无用的后台页面，不仅消耗服务器资源，还会浪费宝贵的爬取预算。", "no_robots_sugg": "在网站根目录创建标准的 robots.txt 文件，并确保其对搜索引擎爬虫公开可见。",
        "robots_bad_rule": "Robots.txt 封禁风险", "robots_bad_rule_desc": "检测到全站封禁规则 (Disallow: /)，且未发现针对 Googlebot 的例外规则。", "robots_bad_rule_impact": "这将直接导致搜索引擎停止抓取并索引您的网站，所有自然搜索流量将归零。", "robots_bad_rule_sugg": "立即移除 'Disallow: /' 规则，或者为搜索引擎爬虫添加具体的 'Allow' 规则。",
        "robots_quality_issue": "Robots.txt 规则配置不当", "robots_quality_issue_desc": "Robots.txt 文件存在潜在问题：{}。", "robots_quality_issue_impact": "可能导致Googlebot行为异常（如误判屏蔽或渲染失败）。", "robots_quality_issue_sugg": "检查 Robots.txt，移除废弃指令（如 Noindex），并确保允许访问 CSS/JS 资源。",
        "robots_no_sitemap": "Robots 未声明 Sitemap", "robots_no_sitemap_desc": "robots.txt 文件中未指明 Sitemap XML 文件的位置。", "robots_no_sitemap_impact": "会降低搜索引擎发现新页面和更新旧内容的速度，尤其对于大型网站影响更明显。", "robots_no_sitemap_sugg": "在 robots.txt 文件底部添加一行：Sitemap: https://yourdomain.com/sitemap.xml",
        "no_sitemap": "Sitemap Failed", "no_sitemap_desc": "Unable to access the Sitemap file (403 Forbidden or 404 Not Found).", "no_sitemap_impact": "Search engines will struggle to find deep links or orphan pages, leading to poor indexing coverage.", "no_sitemap_sugg": "Verify the Sitemap URL is correct and that server permissions allow external access.",
        "sitemap_invalid": "Sitemap 格式错误", "sitemap_invalid_desc": "XML 解析失败，文件格式不符合标准协议。", "sitemap_invalid_impact": "搜索引擎无法读取其中的链接，导致 Sitemap 完全失效。", "sitemap_invalid_sugg": "使用 XML 验证工具检查文件语法，确保没有未闭合的标签或非法字符。",
        "no_favicon": "缺失 Favicon", "no_favicon_desc": "No Favicon icon was detected on the homepage.", "no_favicon_impact": "降低品牌在浏览器标签页和搜索结果页 (SERP) 中的辨识度，进而导致点击率 (CTR) 下降。", "no_favicon_sugg": "制作一个 .ico 或 .png 格式的图标，并在 <head> 中通过 <link rel='icon'> 引用。",
        "duplicate": "发现未规范化的重复内容", "duplicate_desc": "检测到高度相似的内容页面，且未正确配置 Canonical 标签。", "duplicate_impact": "导致关键词内部竞争 (Cannibalization)，分散页面权重，使所有相关页面都难以获得高排名。", "duplicate_sugg": "保留一个首选 URL，并在其他副本页面上添加 rel='canonical' 指向该首选 URL。",
        "http_3xx": "内部链接重定向 (3xx)", "http_3xx_desc": "内部链接发生跳转 (链条: {})。", "http_3xx_impact": "浪费爬虫预算，增加页面加载延迟，且每次跳转都会损耗少量链接传递的权重 (Link Equity)。", "http_3xx_sugg": "批量更新内部链接，使其直接指向最终的目标 URL，避免中间跳转。",
        "http_4xx": "死链/客户端错误 (4xx)", "http_4xx_desc": "内部链接返回 404 (未找到) 或 403 (禁止访问) 错误。", "http_4xx_impact": "严重破坏用户体验，中断权重传递路径，并可能导致已索引的页面被 Google 移除。", "http_4xx_sugg": "移除死链，或者将其重定向到最相关的有效页面。",
        "http_5xx": "服务器错误 (5xx)", "http_5xx_desc": "服务器响应 500/502/503 等内部错误。", "http_5xx_impact": "表明服务器极其不稳定，Googlebot 会因此降低对该站点的爬取频率以减轻负载。", "http_5xx_sugg": "检查服务器错误日志，优化数据库查询或升级服务器配置。",
        "hreflang_invalid": "Hreflang 格式错误", "hreflang_invalid_desc": "语言代码不符合 ISO 639-1 标准 (如使用了 {} 等错误格式)。", "hreflang_invalid_impact": "Google 无法识别目标语言，导致国际化定位失效。", "hreflang_invalid_sugg": "使用标准的 ISO 语言代码 (例如 'en-US' 而不是 'en_US')。",
        "hreflang_no_default": "Hreflang 缺失 x-default", "hreflang_no_default_desc": "Missing 'x-default' fallback attribute.", "hreflang_no_default_impact": "当用户来自未指定的语言/地区时，可能无法自动匹配到最合适的通用版本（通常是英语）。", "hreflang_no_default_sugg": "添加 hreflang='x-default' 标签，指定默认的语言版本。",
        "alt_bad_quality": "图片 Alt 质量差", "alt_bad_quality_desc": "Alt 文本使用了无意义词汇（如 image1.jpg, photo）或过短。", "alt_bad_quality_impact": "搜索引擎无法理解图片内容，错失图片搜索流量，且对视障用户极不友好。", "alt_bad_quality_sugg": "使用描述性文本准确描述图片内容，包含相关的关键词。",
        "anchor_bad_quality": "锚文本质量差", "anchor_bad_quality_desc": "使用了“点击这里”、“更多”等通用词汇作为链接文本。", "anchor_bad_quality_impact": "无法向搜索引擎传递目标页面的关键词相关性，降低了目标页面的排名潜力。", "anchor_bad_quality_sugg": "使用描述性 keywords in the anchor text.",
        "cls_risk": "CLS 布局偏移风险 (静态检测)", "cls_risk_desc": "检测到 <img> 标签缺失 width 或 height 属性。", "cls_risk_impact": "图片加载时会撑开页面，导致布局发生意外抖动，直接恶化 CLS 指标。", "cls_risk_sugg": "在 HTML 中显式指定图片和视频的宽度和高度属性。",
        "missing_title": "缺失页面标题 (Title)", "missing_title_desc": "页面代码中未找到 <title> 标签。", "missing_title_impact": "Title 是最重要的 SEO 标签。缺失将导致搜索引擎无法判断页面主题，关键词排名极差。", "missing_title_sugg": "为每个页面添加独特、包含核心关键词的标题。",
        "short_title": "标题过短", "short_title_desc": "标题长度不足 (约 {} px)，难以完整表达页面意图。", "short_title_impact": "浪费了宝贵的标题空间，错失了覆盖长尾关键词排名的机会。", "short_title_sugg": "丰富标题内容，加入品牌词或修饰词，建议长度在 285-575 px 之间。",
        "long_title": "标题过长", "long_title_desc": "标题超过建议显示宽度 (约 {} px)。", "long_title_impact": "标题将在搜索结果中被截断，降低可读性和点击率。", "long_title_sugg": "精简标题长度，将核心信息前置，控制在 600 px 以内。",
        "missing_desc": "缺失元描述", "missing_desc_desc": "页面未包含 <meta name='description'> 标签。", "missing_desc_impact": "Google will generate a snippet from page text, which is often irrelevant and lowers CTR.", "missing_desc_sugg": "添加吸引人的元描述，概括页面内容并包含号召性用语。",
        "short_desc": "元描述过短", "short_desc_desc": "内容过少 (约 {} px)，吸引力不足。", "short_desc_impact": "无法充分展示页面卖点，在搜索结果中缺乏竞争力。", "short_desc_sugg": "扩充描述至 400-920 px，提供更多有价值的信息。",
        "missing_h1": "缺失 H1 标签", "missing_h1_desc": "页面缺乏 <h1> 主标题。", "missing_h1_impact": "搜索引擎难以理解内容的层级结构和核心主题，降低了关键词的相关性权重。", "missing_h1_sugg": "确保每个页面有且仅有一个 H1 标签，概括当前页面的主题。",
        "missing_viewport": "缺失移动端视口配置", "missing_viewport_desc": "未配置 <meta name='viewport'> 标签。", "missing_viewport_impact": "在移动设备上显示异常（字体极小）。Google 移动优先索引会严重惩罚此类页面。", "missing_viewport_sugg": "在 <head> 中添加标准的 viewport meta 标签。",
        "missing_canonical": "缺失 Canonical 标签", "missing_canonical_desc": "未指定规范链接。", "missing_canonical_impact": "无法应对 URL 参数（如 ?id=1）导致的重复内容问题，容易造成权重稀释。", "missing_canonical_sugg": "在所有页面添加自引用（Self-referencing）或指向原件的 Canonical 标签。",
        "missing_jsonld": "缺失结构化数据", "missing_jsonld_desc": "未检测到 Schema.org 标记。", "missing_jsonld_impact": "错失富媒体搜索结果（Rich Results），在 SERP 中不如竞争对手显眼。", "missing_jsonld_sugg": "建议配置结构化数据。基于页面内容，推荐添加：{}。",
        "missing_hreflang": "缺失 Hreflang", "missing_hreflang_desc": "未发现语言区域标记（HTML或Sitemap中均无）。", "missing_hreflang_impact": "多语言站点无法正确定位目标受众，导致流量不精准。", "missing_hreflang_sugg": "在 HTML 头部或 Sitemap 中配置 hreflang 标签。",
        "soft_404": "疑似软 404 (Soft 404)", "soft_404_desc": "页面返回 200 状态码但内容显示“未找到”。", "soft_404_impact": "严重浪费爬虫预算，导致无效页面挤占有效页面的索引名额。", "soft_404_sugg": "配置服务器对不存在的页面返回 404 HTTP 状态码。",
        "missing_alt": "图片缺失 Alt 属性", "missing_alt_desc": "图片标签缺少 alt 属性。", "missing_alt_impact": "搜索引擎无法理解图片内容，错失图片搜索流量。", "missing_alt_sugg": "为所有有意义的图片添加描述性的 alt 属性。",
        "js_links": "发现 JS 伪链接", "js_links_desc": "使用了 href='javascript:...' 形式的链接。", "js_links_impact": "爬虫无法跟踪此类链接，导致内部链接断裂，深层页面变成“孤岛”。", "js_links_sugg": "使用标准的 <a href> 标签，仅在 onclick 事件中处理 JS 逻辑。",
        "url_underscore": "URL 包含下划线", "url_underscore_desc": "URL 路径中使用下划线 (_) 分隔单词。", "url_underscore_impact": "Google 建议使用连字符。下划线可能导致关键词无法被正确切分（被视为一个长单词）。", "url_underscore_sugg": "在 URL 结构中使用连字符 (-) 代替下划线。",
        "url_uppercase": "URL 包含大写字母", "url_uppercase_desc": "URL 路径中混用了大写字母。", "url_uppercase_impact": "服务器通常区分大小写，极易造成一页多址（Duplicate Content）和 404 错误。", "url_uppercase_sugg": "强制所有 URL 使用小写字母。",
        
        # Baidu specific
        "missing_keywords": "Missing Meta Keywords (Baidu)",
        "missing_keywords_desc": "No <meta name='keywords'> tag found.",
        "missing_keywords_impact": "Baidu still uses keywords as a ranking signal, unlike Google.",
        "missing_keywords_sugg": "Add meta keywords tag with 3-5 relevant keywords.",
        "missing_baidu_stats": "Missing Baidu Analytics",
        "missing_baidu_stats_desc": "Baidu Tongji script (hm.baidu.com) not found.",
        "missing_baidu_stats_impact": "Unable to track Baidu traffic effectively.",
        "missing_baidu_stats_sugg": "Install Baidu Tongji script.",
        "missing_baidu_verify": "Missing Baidu Verification",
        "missing_baidu_verify_desc": "No 'baidu-site-verification' tag found.",
        "missing_baidu_verify_impact": "May delay site indexing on Baidu.",
        "missing_baidu_verify_sugg": "Add verification tag.",
        "baidu_robots_missing": "Missing Baidu Rules",
        "baidu_robots_missing_desc": "No specific rules for 'Baiduspider' in Robots.txt.",
        "baidu_robots_missing_impact": "Inefficient crawling by Baidu.",
        "baidu_robots_missing_sugg": "Add User-agent: Baiduspider directives.",
        "missing_applicable_device": "Missing Applicable Device (Baidu)",
        "missing_applicable_device_desc": "Meta tag 'applicable-device' not found.",
        "missing_applicable_device_impact": "Baidu can't identify if page is PC/Mobile adapted.",
        "missing_applicable_device_sugg": "Add <meta name='applicable-device' content='pc,mobile'>.",
        "missing_no_transform": "Missing No-transform (Baidu)",
        "missing_no_transform_desc": "Cache-Control: no-transform not found.",
        "missing_no_transform_impact": "Baidu might transcode your page (Siteapp), breaking layout.",
        "missing_no_transform_sugg": "Add <meta http-equiv='Cache-Control' content='no-transform'>.",
        "missing_icp": "Missing ICP Number",
        "missing_icp_desc": "No ICP filing number found in page content.",
        "missing_icp_impact": "Required by Chinese law for mainland hosting; affects trust and Baidu ranking.",
        "missing_icp_sugg": "Add ICP number in footer linking to beian.miit.gov.cn.",
        "content_not_chinese": "Low Chinese Content Ratio",
        "content_not_chinese_desc": "Chinese character ratio is below 5%.",
        "content_not_chinese_impact": "Baidu prioritizes Chinese content. Low ratio affects ranking in CN search.",
        "content_not_chinese_sugg": "Ensure main content is in Simplified Chinese.",
        "server_not_in_china": "Server Not In China (Baidu)",
        "server_not_in_china_desc": "Server IP detected in: {}. Baidu prefers mainland China hosting.",
        "server_not_in_china_impact": "Slow cross-border loading may cause Baidu spider timeouts.",
        "server_not_in_china_sugg": "Migrate hosting to Mainland China and get ICP filing."
    },
    "en": {
        "sidebar_title": "🔍 AuditAI Pro",
        "sidebar_caption": "Deep Audit Edition v14.2",
        "nav_label": "Navigation",
        "nav_options": ["Input URL", "Dashboard", "Data Matrix", "PPT Generator"],
        "lang_label": "Language / 语言",
        "clear_data": "Clear Data & Reset",
        "cache_info": "Cached {} pages",
        "sitemap_status_title": "Sitemap Status:",
        "sitemap_found_href": "✅ Hreflang Found", 
        "sitemap_no_href": "⚠️ No Hreflang",       
        "sitemap_missing": "❌ Sitemap Missing",
        
        # PSI Related
        "psi_settings": "Google PSI API Settings (Optional)",
        "psi_api_key_label": "Enter Google PageSpeed API Key",
        "psi_api_help": "Enter API Key to fetch Real User Metrics (LCP, CLS, INP) for the home page. Leave empty for code-only check.",
        "psi_list_url_label": "Product List URL (Optional)", 
        "psi_detail_url_label": "Product Detail URL (Optional)", 
        "psi_get_key": "No API Key? [Get one for free here](https://developers.google.com/speed/docs/insights/v5/get-started)",
        "psi_fetching": "Fetching real CWV data from Google API ({}) ...",
        "psi_success": "Real user data fetched successfully!",
        "psi_error": "API Failed or No CrUX Data",
        
        "input_header": "Start Deep Audit",
        "input_info": "Note: v14.2 features Manual URL input.",
        "input_label": "Target URL",
        "input_placeholder": "https://example.com",
        "max_pages_label": "Max Pages to Crawl",
        "adv_settings": "Advanced Settings", 
        "check_robots_label": "Check & Respect Robots.txt", 
        "crawl_sitemap_label": "Parse Sitemap from Robots.txt", 
        "baidu_mode_label": "Enable Baidu SEO Audit Mode", 
        "allow_subdomains_label": "Allow Subdomains (e.g. blog.site.com)", 
        "allow_outside_folder_label": "Allow Outside Start Folder (e.g. /fr/ from /en/)", 
        "manual_sitemaps": "Manual Sitemap URLs (One per line, Optional)", 
        "manual_pages_label": "Manual Pages to Audit (One per line)",
        "sitemap_content_label": "Paste Sitemap XML Content (Direct Parse)",
        "start_btn": "Start Deep Crawl",
        "error_url": "Invalid URL format",
        "spinner_crawl": "Running Deep Audit (Max {} pages)...", 
        "error_no_data": "No pages crawled. Reason: {}", 
        "success_audit": "Audit Complete! Analyzed {} pages.",
        
        "dashboard_header": "Executive Summary",
        "warn_no_data": "No data available.",
        "kpi_health": "Health Score",
        "kpi_pages": "Analyzed Pages",
        "kpi_issues": "Total Issues",
        "kpi_critical": "Critical Issues",
        "chart_issues": "Issue Distribution",
        "chart_no_issues": "No significant issues found.",
        "chart_status": "HTTP Status Codes",
        "cwv_title": "Core Web Vitals - Real User Data (Home Only)",
        "cwv_source": "Source: Google Chrome User Experience Report (CrUX)",
        
        "matrix_header": "Crawled Data Matrix",
        "download_csv": "Download CSV Report",
        
        "ppt_header": "Pitch Deck Preview",
        "ppt_success_no_issues": "No critical issues found.",
        "ppt_download_header": "📥 Export Report",
        "ppt_info": "Note: PPT optimized for 16:9 with logical issue ordering.",
        "ppt_btn": "Generate & Download .pptx",
        "ppt_preview_header": "Web Preview",
        "ppt_slide_title": "Issue Type:",
        "ppt_category": "Category:",
        "ppt_severity": "Severity:",
        "ppt_impact": "Impact:",
        "ppt_impact_desc": "Affects **{}** pages in crawled sample.",
        "ppt_desc": "Description:",
        "ppt_business_impact": "📉 Business & SEO Impact:",
        "ppt_sugg": "💡 Suggestion:",
        "ppt_examples": "🔍 Examples:",
        "ppt_prev": "⬅️ Previous",
        "ppt_next": "Next ➡️",
        
        # Categories
        "cat_access": "Access & Indexing",
        "cat_indexability": "Indexability",
        "cat_technical": "Technical SEO",
        "cat_content": "On-Page Content",
        "cat_image_ux": "UX & Assets",
        "cat_cwv_performance": "Core Web Vitals (Performance)",
        
        "ppt_cover_title": "SEO Technical Audit",
        "ppt_cover_sub": "Generated by AuditAI Pro v14.0",
        "ppt_slide_desc_title": "Description & Impact",
        "ppt_slide_count_title": "Affected Pages (in sample): {}",
        "ppt_slide_ex_title": "Affected Page Examples",
        "ppt_slide_sugg_title": "💡 Recommendation:",
        "serp_sim_title": "Google SERP Simulation:",
        "rich_sim_title": "Rich Results Simulation:",
        "code_sim_title": "Code Snippet:",
        "visual_sim_title": "Visual Experience Simulation:",
        "cwv_sim_title": "CWV Metric Visualization:",

        # Issues
        "lcp_issue": "LCP (Largest Contentful Paint) Fails", 
        "lcp_issue_desc": "LCP is {:.2f}s (Target <2.5s). Main content takes too long to appear.", 
        "lcp_issue_impact": "LCP is a core ranking factor. Slow loading speeds significantly increase bounce rates and lower search rankings.", 
        "lcp_issue_sugg": "Optimize image sizes (use WebP), implement a CDN, defer non-critical JavaScript, and preload the LCP element.",
        
        "cls_issue": "CLS (Cumulative Layout Shift) Fails", 
        "cls_issue_desc": "CLS score is {:.3f} (Target <0.1). Elements on the page shift unexpectedly during loading.", 
        "cls_issue_impact": "A Core Web Vital ranking factor. Visual instability frustrates users and can lead to accidental clicks, damaging brand reputation.", 
        "cls_issue_sugg": "Set explicit width and height attributes for all images and videos, and avoid inserting dynamic content above the fold.",
        
        "inp_issue": "INP (Interaction to Next Paint) Fails", 
        "inp_issue_desc": "INP is {}ms (Target <200ms). The page is unresponsive to user clicks or interactions.", 
        "inp_issue_impact": "A new Core Web Vital. High latency makes the site feel 'broken' or sluggish, severely impacting user conversion rates.", 
        "inp_issue_sugg": "Reduce main-thread blocking, break up Long Tasks, and optimize complex JavaScript event handlers.",
        "fcp_issue": "FCP (First Contentful Paint) Slow", "fcp_issue_desc": "FCP is {:.2f}s (Target <1.8s). Time to see first content is too long.", "fcp_issue_impact": "Slow FCP increases bounce rate as users feel the site is unresponsive.", "fcp_issue_sugg": "Improve TTFB and eliminate render-blocking resources.",
        "fcp_baidu_issue": "FCP not meeting Baidu standards", "fcp_baidu_issue_desc": "FCP is {:.2f}s (Baidu Target < 2.0s).", "fcp_baidu_issue_impact": "Baidu has strict requirements for mobile landing page speed. >2s leads to ranking penalty.", "fcp_baidu_issue_sugg": "Use China CDN, optimize critical path, ensure FCP < 2s.",
        
        "no_robots": "Missing Robots.txt", 
        "no_robots_desc": "The robots.txt file was not found in the root directory, or the server returned an error.", 
        "no_robots_impact": "Search engines may index useless or admin pages, wasting your crawl budget and server resources.", 
        "no_robots_sugg": "Create a standard robots.txt file in the root directory and ensure it is publicly accessible.",
        
        "robots_bad_rule": "Robots.txt Blocking", 
        "robots_bad_rule_desc": "A global blocking rule (Disallow: /) was detected.", 
        "robots_bad_rule_impact": "This prevents search engines from crawling your entire site, resulting in total de-indexing and zero organic traffic.", 
        "robots_bad_rule_sugg": "Remove the 'Disallow: /' rule immediately to allow crawling.",
        
        "robots_quality_issue": "Robots.txt Configuration", 
        "robots_quality_issue_desc": "Potential issue in Robots.txt: {}.", 
        "robots_quality_issue_impact": "May block important resources like CSS/JS, affecting how Google renders the page.", 
        "robots_quality_issue_sugg": "Review Robots.txt to ensure CSS and JS files are crawlable.",
        
        "robots_no_sitemap": "Sitemap Missing in Robots", 
        "robots_no_sitemap_desc": "The location of the Sitemap XML is not specified in the robots.txt file.", 
        "robots_no_sitemap_impact": "This slows down the discovery of new pages and content updates, especially for larger websites.", 
        "robots_no_sitemap_sugg": "Add a 'Sitemap: [URL]' directive to the bottom of your robots.txt file.",
        
        "no_sitemap": "Sitemap Failed", 
        "no_sitemap_desc": "Unable to access the Sitemap file (403 Forbidden or 404 Not Found).", 
        "no_sitemap_impact": "Search engines will struggle to find deep links or orphan pages, leading to poor indexing coverage.", 
        "no_sitemap_sugg": "Verify the Sitemap URL is correct and that server permissions allow external access.",
        
        "sitemap_invalid": "Invalid Sitemap", 
        "sitemap_invalid_desc": "XML parsing failed. The file format does not adhere to the standard protocol.", 
        "sitemap_invalid_impact": "Search engines cannot read the links inside, rendering the Sitemap completely useless.", 
        "sitemap_invalid_sugg": "Validate your XML syntax to ensure there are no unclosed tags or invalid characters.",
        
        "no_favicon": "Missing Favicon", 
        "no_favicon_desc": "No Favicon icon was detected on the homepage.", 
        "no_favicon_impact": "Reduces brand visibility in browser tabs and lowers the Click-Through Rate (CTR) in search results.", 
        "no_favicon_sugg": "Create a .ico or .png icon and link it in the <head> section.",
        
        "duplicate": "Duplicate Content", 
        "duplicate_desc": "Identical content detected across multiple URLs without proper canonicalization.", 
        "duplicate_impact": "Causes keyword cannibalization and dilutes link equity, preventing both pages from ranking well.", 
        "duplicate_sugg": "Select a preferred URL and use rel='canonical' tags on duplicate versions to point to it.",
        
        "http_3xx": "Redirect Chain", 
        "http_3xx_desc": "Internal link triggers a redirect (Chain: {}).", 
        "http_3xx_impact": "Wastes crawl budget, adds latency to page load, and dilutes the link equity passed to the destination.", 
        "http_3xx_sugg": "Update the internal link to point directly to the final destination URL.",
        
        "http_4xx": "Broken Link", 
        "http_4xx_desc": "Internal link returns a 4xx error (e.g., 404 Not Found).", 
        "http_4xx_impact": "Creates a bad user experience, breaks the flow of link equity, and may cause indexed pages to be dropped.", 
        "http_4xx_sugg": "Fix the broken link or remove it.",
        
        "http_5xx": "Server Error", 
        "http_5xx_desc": "Server returned a 5xx error (e.g., 500 Internal Server Error).", 
        "http_5xx_impact": "Signals server instability. Googlebot will reduce the crawl rate of your site to avoid overloading it.", 
        "http_5xx_sugg": "Check server error logs and ensure database stability.",
        
        "hreflang_invalid": "Invalid Hreflang", 
        "hreflang_invalid_desc": "The language code format does not comply with ISO 639-1 standards.", 
        "hreflang_invalid_impact": "Google cannot identify the target language, causing international targeting to fail.", 
        "hreflang_invalid_sugg": "Use standard ISO codes (e.g., 'en-US' instead of 'en_US').",
        
        "hreflang_no_default": "No x-default", 
        "hreflang_no_default_desc": "Missing 'x-default' fallback attribute.", 
        "hreflang_no_default_impact": "Users from unspecified regions may be served the wrong language version.", 
        "hreflang_no_default_sugg": "Add an hreflang='x-default' tag to specify the default version.",
        
        "alt_bad_quality": "Bad Alt Text", 
        "alt_bad_quality_desc": "Alt text uses filenames or generic words like 'image'.", 
        "alt_bad_quality_impact": "Search engines cannot understand the image context, hurting Image SEO and accessibility.", 
        "alt_bad_quality_sugg": "Use descriptive text that accurately describes the image content.",
        
        "anchor_bad_quality": "Bad Anchor", 
        "anchor_bad_quality_desc": "Generic anchor text found (e.g., 'Click here').", 
        "anchor_bad_quality_impact": "Fails to pass keyword relevance to the target page, reducing its ranking potential.", 
        "anchor_bad_quality_sugg": "Use descriptive keywords in the anchor text.",
        
        "cls_risk": "CLS Risk (Static)", 
        "cls_risk_desc": "Images missing width or height attributes detected.", 
        "cls_risk_impact": "Images will push content down as they load, causing layout shifts and hurting Core Web Vitals.", 
        "cls_risk_sugg": "Explicitly set width and height attributes on all image tags.",
        
        "missing_title": "Missing Title", 
        "missing_title_desc": "No <title> tag found in the page code.", 
        "missing_title_impact": "Title is the most important on-page SEO factor. Missing it causes severe ranking loss.", 
        "missing_title_sugg": "Add a unique, keyword-rich title to every page.",
        
        "short_title": "Title Short", 
        "short_title_desc": "Title is too short (~{}px).", 
        "short_title_impact": "Missed opportunity to target relevant keywords and attract clicks.", 
        "short_title_sugg": "Expand the title to ~285-575px, including your brand name.",
        
        "long_title": "Title Long", 
        "long_title_desc": "Title exceeds optimal width (~{}px).", 
        "long_title_impact": "The title will be truncated in search results, reducing readability and CTR.", 
        "long_title_sugg": "Shorten the title to under 600px, keeping important keywords at the front.",
        
        "missing_desc": "Missing Description", 
        "missing_desc_desc": "No meta description tag found.", 
        "missing_desc_impact": "Google will generate a snippet from page text, which is often irrelevant and lowers CTR.", 
        "missing_desc_sugg": "Add a compelling meta description that summarizes the page content.",
        
        "short_desc": "Description Short", 
        "short_desc_desc": "Description content is too thin (~{}px).", 
        "short_desc_impact": "Fails to provide enough context to entice users to click.", 
        "short_desc_sugg": "Expand the description to 400-920px with a call to action.",
        
        "missing_h1": "Missing H1", 
        "missing_h1_desc": "No <h1> heading tag found.", 
        "missing_h1_impact": "Search engines struggle to identify the main topic of the page.", 
        "missing_h1_sugg": "Ensure every page has exactly one H1 tag describing the main topic.",
        
        "missing_viewport": "No Viewport", 
        "missing_viewport_desc": "Mobile viewport meta tag is missing.", 
        "missing_viewport_impact": "The page is not mobile-friendly. Google Mobile-First Indexing will severely penalize it.", 
        "missing_viewport_sugg": "Add the standard viewport meta tag to the <head>.",
        
        "missing_canonical": "No Canonical", 
        "missing_canonical_desc": "Missing canonical tag.", 
        "missing_canonical_impact": "High risk of duplicate content issues, especially with URL parameters.", 
        "missing_canonical_sugg": "Add a self-referencing canonical tag to all pages.",
        
        "missing_jsonld": "No Schema", 
        "missing_jsonld_desc": "No JSON-LD structured data found.", 
        "missing_jsonld_impact": "Missed opportunity for Rich Snippets (e.g., Stars, Price) which boost CTR.", 
        "missing_jsonld_sugg": "Add JSON-LD schema based on page type: {}.",
        
        "missing_hreflang": "No Hreflang", 
        "missing_hreflang_desc": "No language targeting tags found.", 
        "missing_hreflang_impact": "Poor international targeting.", 
        "missing_hreflang_sugg": "Add hreflang tags.",
        
        "soft_404": "Soft 404", 
        "soft_404_desc": "Page returns a 200 OK status but displays an error message.", 
        "soft_404_impact": "Wastes crawl budget on invalid pages and confuses search engines.", 
        "soft_404_sugg": "Configure the server to return a 404 Not Found status code.",
        
        "missing_alt": "Missing Alt", 
        "missing_alt_desc": "Images lack alternative text attributes.", 
        "missing_alt_impact": "Bad for accessibility and prevents images from ranking in Image Search.", 
        "missing_alt_sugg": "Add descriptive alt text to all relevant images.",
        
        "js_links": "JS Links", 
        "js_links_desc": "Uncrawlable JavaScript links found.", 
        "js_links_impact": "Search engines cannot follow these links, leaving pages orphaned.", 
        "js_links_sugg": "Replace with standard <a href> tags.",
        
        "url_underscore": "URL Underscores", 
        "url_underscore_desc": "URL uses underscores to separate words.", 
        "url_underscore_impact": "Google treats underscores as joiners, not separators, hurting keyword parsing.", 
        "url_underscore_sugg": "Use hyphens (-) instead of underscores.",
        
        "url_uppercase": "URL Uppercase", 
        "url_uppercase_desc": "URL contains uppercase letters.", 
        "url_uppercase_impact": "Can lead to duplicate content issues on case-sensitive servers.", 
        "url_uppercase_sugg": "Force all URLs to be lowercase.",
        
        # Baidu
        "missing_keywords": "Missing Meta Keywords (Baidu)",
        "missing_keywords_desc": "No <meta name='keywords'> tag found.",
        "missing_keywords_impact": "Baidu still uses keywords as a ranking signal, unlike Google.",
        "missing_keywords_sugg": "Add meta keywords tag with 3-5 relevant keywords.",
        "missing_baidu_stats": "Missing Baidu Analytics",
        "missing_baidu_stats_desc": "Baidu Tongji script (hm.baidu.com) not found.",
        "missing_baidu_stats_impact": "Unable to track Baidu traffic effectively.",
        "missing_baidu_stats_sugg": "Install Baidu Tongji script.",
        "missing_baidu_verify": "Missing Baidu Verification",
        "missing_baidu_verify_desc": "No 'baidu-site-verification' tag found.",
        "missing_baidu_verify_impact": "May delay site indexing on Baidu.",
        "missing_baidu_verify_sugg": "Add verification tag.",
        "baidu_robots_missing": "Missing Baidu Rules",
        "baidu_robots_missing_desc": "No specific rules for 'Baiduspider' in Robots.txt.",
        "baidu_robots_missing_impact": "Inefficient crawling by Baidu.",
        "baidu_robots_missing_sugg": "Add User-agent: Baiduspider directives.",
        "missing_applicable_device": "Missing Applicable Device (Baidu)",
        "missing_applicable_device_desc": "Meta tag 'applicable-device' not found.",
        "missing_applicable_device_impact": "Baidu can't identify if page is PC/Mobile adapted.",
        "missing_applicable_device_sugg": "Add <meta name='applicable-device' content='pc,mobile'>.",
        "missing_no_transform": "Missing No-transform (Baidu)",
        "missing_no_transform_desc": "Cache-Control: no-transform not found.",
        "missing_no_transform_impact": "Baidu might transcode your page (Siteapp), breaking layout.",
        "missing_no_transform_sugg": "Add <meta http-equiv='Cache-Control' content='no-transform'>.",
        "missing_icp": "Missing ICP Number",
        "missing_icp_desc": "No ICP filing number found in page content.",
        "missing_icp_impact": "Required by Chinese law for mainland hosting; affects trust and Baidu ranking.",
        "missing_icp_sugg": "Add ICP number in footer linking to beian.miit.gov.cn.",
        "content_not_chinese": "Low Chinese Content Ratio",
        "content_not_chinese_desc": "Chinese character ratio is below 5%.",
        "content_not_chinese_impact": "Baidu prioritizes Chinese content. Low ratio affects ranking in CN search.",
        "content_not_chinese_sugg": "Ensure main content is in Simplified Chinese.",
        "server_not_in_china": "Server Not In China (Baidu)",
        "server_not_in_china_desc": "Server IP detected in: {}. Baidu prefers mainland China hosting.",
        "server_not_in_china_impact": "Slow cross-border loading may cause Baidu spider timeouts.",
        "server_not_in_china_sugg": "Migrate hosting to Mainland China and get ICP filing."
    }
}

# --- Level 6: 核心逻辑 (Data Layer) ---
def get_translated_text(issue_id, lang, args=None):
    if args is None: args = []
    t = TRANSLATIONS[lang]
    
    def safe_format(text, arguments):
        try: return text.format(*arguments)
        except IndexError: return text

    return {
        "title": t.get(issue_id, issue_id),
        "desc": safe_format(t.get(issue_id + "_desc", ""), args),
        "impact": t.get(issue_id + "_impact", ""),
        "suggestion": safe_format(t.get(issue_id + "_sugg", ""), args)
    }

def fetch_psi_data(url, api_key):
    if not api_key: return None
    endpoint = f"https://www.googleapis.com/pagespeedonline/v5/runPagespeed?url={url}&key={api_key}&strategy=mobile"
    try:
        response = requests.get(endpoint, timeout=30)
        if response.status_code == 200:
            data = response.json()
            crux = data.get('loadingExperience', {}).get('metrics', {})
            if not crux: return {"error": "No CrUX data available"}
            return {
                "LCP": crux.get('LARGEST_CONTENTFUL_PAINT_MS', {}).get('percentile', 0) / 1000,
                "CLS": crux.get('CUMULATIVE_LAYOUT_SHIFT_SCORE', {}).get('percentile', 0) / 100,
                "INP": crux.get('INTERACTION_TO_NEXT_PAINT', {}).get('percentile', 0),
                "FCP": crux.get('FIRST_CONTENTFUL_PAINT_MS', {}).get('percentile', 0) / 1000,
            }
        else: return {"error": f"API Error: {response.status_code}"}
    except Exception as e: return {"error": str(e)}

def check_cwv_issues(cwv_data, url, label=""):
    issues = []
    if not cwv_data or "error" in cwv_data: return issues
    category_key = "cwv_performance"
    
    lcp = cwv_data.get("LCP", 0)
    if lcp > 2.5:
        issues.append({
            "id": "lcp_issue", "category": category_key, "severity": "Critical" if lcp > 4.0 else "High",
            "url": url, "args": [lcp], "examples": [f"{url} ({lcp:.2f}s) {label}"] 
        })
    
    inp = cwv_data.get("INP", 0)
    if inp > 200:
        issues.append({
            "id": "inp_issue", "category": category_key, "severity": "Critical" if inp > 500 else "High",
            "url": url, "args": [inp], "examples": [f"{url} ({inp}ms) {label}"]
        })

    cls = cwv_data.get("CLS", 0)
    if cls > 0.1:
        issues.append({
            "id": "cls_issue", "category": category_key, "severity": "Critical" if cls > 0.25 else "High",
            "url": url, "args": [cls], "examples": [f"{url} ({cls:.3f}) {label}"]
        })
    
    fcp = cwv_data.get("FCP", 0)
    if fcp > 1.8:
        issues.append({
            "id": "fcp_issue", "category": category_key, "severity": "Critical" if fcp > 3.0 else "Medium",
            "url": url, "args": [fcp], "examples": [f"{url} ({fcp:.2f}s) {label}"]
        })

    return issues

def check_site_level_assets(start_url, lang="zh", check_robots=True, crawl_sitemap_flag=True, manual_sitemaps=None, baidu_mode=False):
    issues = []
    sitemap_has_hreflang = False
    
    initial_netloc = urlparse(start_url).netloc
    base_url = f"{urlparse(start_url).scheme}://{initial_netloc}"
    headers = get_browser_headers()
    
    robots_url = urljoin(base_url, "/robots.txt")
    if check_robots:
        try:
            r = requests.get(robots_url, headers=headers, timeout=10, allow_redirects=True, stream=True, verify=False)
            if r.status_code != 200:
                issues.append({"id": "no_robots", "category": "access", "severity": "Medium", "url": robots_url, "examples": [robots_url]})
            else:
                content = r.text.lower()
                if len(content.strip()) < 5:
                     issues.append({"id": "robots_quality_issue", "category": "access", "severity": "Medium", "url": robots_url, "args": ["File is empty or too short"], "examples": [robots_url]})
                if "user-agent" not in content:
                     issues.append({"id": "robots_quality_issue", "category": "access", "severity": "Medium", "url": robots_url, "args": ["Missing User-agent directive"], "examples": [robots_url]})
                if "disallow: /*.css" in content or "disallow: /*.js" in content:
                     issues.append({"id": "robots_quality_issue", "category": "access", "severity": "High", "url": robots_url, "args": ["Blocking CSS/JS resources"], "examples": [robots_url]})
                
                if "disallow: /" in content and "allow:" not in content:
                    issues.append({"id": "robots_bad_rule", "category": "access", "severity": "Critical", "url": robots_url, "examples": [robots_url]})
                
                if baidu_mode:
                    if "baiduspider" not in content:
                        issues.append({"id": "baidu_robots_missing", "category": "access", "severity": "Low", "url": robots_url, "examples": [robots_url]})
                
                if "sitemap:" not in content:
                    issues.append({"id": "robots_no_sitemap", "category": "access", "severity": "Low", "url": robots_url, "examples": [robots_url]})
                
                if crawl_sitemap_flag:
                    sitemaps_in_robots = re.findall(r'sitemap:\s*(https?://\S+)', content, re.IGNORECASE)
                    if sitemaps_in_robots:
                        if manual_sitemaps is None: manual_sitemaps = []
                        manual_sitemaps.extend(sitemaps_in_robots)
            r.close()
        except: 
            issues.append({"id": "no_robots", "category": "access", "severity": "Medium", "url": robots_url, "examples": [robots_url]})

    sitemap_urls = manual_sitemaps if manual_sitemaps else [urljoin(base_url, "/sitemap.xml")]
    any_valid = False
    for sm_url in sitemap_urls:
        if not sm_url.strip(): continue
        try:
            r = requests.get(sm_url, headers=headers, timeout=15, verify=False)
            if r.status_code == 200:
                try:
                    ET.fromstring(r.content)
                    any_valid = True
                    if 'hreflang' in r.text or 'xhtml' in r.text: sitemap_has_hreflang = True
                except:
                    if not sm_url.endswith('.gz'):
                        issues.append({"id": "sitemap_invalid", "category": "access", "severity": "Medium", "url": sm_url, "examples": [sm_url]})
            else:
                if manual_sitemaps: issues.append({"id": "no_sitemap", "category": "access", "severity": "Low", "url": sm_url, "examples": [sm_url]})
        except:
            if manual_sitemaps: issues.append({"id": "no_sitemap", "category": "access", "severity": "Low", "url": sm_url, "examples": [sm_url]})

    if not any_valid and not manual_sitemaps:
         issues.append({"id": "no_sitemap", "category": "access", "severity": "Low", "url": sitemap_urls[0], "examples": [sitemap_urls[0]]})

    try:
        r = requests.get(urljoin(base_url, "/favicon.ico"), headers=headers, timeout=5, verify=False)
        if r.status_code != 200 or int(r.headers.get('content-length', 0)) == 0:
            issues.append({"id": "no_favicon", "category": "image_ux", "severity": "Low", "url": base_url, "examples": [base_url]})
    except: pass
    
    if baidu_mode:
        cc, country_name = check_server_location(start_url)
        if cc and cc != 'CN':
             issues.append({"id": "server_not_in_china", "category": "technical", "severity": "High", "url": start_url, "args": [country_name], "examples": [start_url]})

    return issues, sitemap_has_hreflang

def analyze_page(url, content, status, sitemap_has_hreflang, baidu_mode=False):
    soup = BeautifulSoup(content, 'html.parser')
    issues = []
    
    title = soup.title.string.strip() if soup.title else None
    desc = soup.find('meta', attrs={'name': 'description'})
    desc_content = desc['content'].strip() if desc else None
    h1 = soup.find('h1')
    h1_content = h1.get_text().strip() if h1 else None
    
    can_tag = soup.find('link', attrs={'rel': 'canonical'})
    can_url = can_tag['href'] if can_tag else None

    if status == 200:
        is_self_canonical = True
        if can_url:
            def norm_u(u): return u.split('#')[0].rstrip('/')
            try:
                abs_can = urljoin(url, can_url)
                if norm_u(abs_can) != norm_u(url):
                    is_self_canonical = False
            except: pass

        if not can_url:
            issues.append({"id": "missing_canonical", "category": "indexability", "severity": "Medium", "url": url})
            is_self_canonical = True

        hreflangs = soup.find_all('link', hreflang=True)
        if hreflangs:
            has_x_default = False
            invalid = []
            pat = re.compile(r'^[a-z]{2}(-[a-zA-Z]{2})?$|x-default', re.IGNORECASE)
            for link in hreflangs:
                code = link.get('hreflang', '').strip()
                if code.lower() == 'x-default': has_x_default = True
                if not pat.match(code): invalid.append(code)
            if invalid:
                issues.append({"id": "hreflang_invalid", "category": "indexability", "severity": "High", "url": url, "args": [", ".join(invalid[:3])]})
            if not has_x_default:
                issues.append({"id": "hreflang_no_default", "category": "indexability", "severity": "Low", "url": url})
        elif not sitemap_has_hreflang:
             if is_self_canonical:
                issues.append({"id": "missing_hreflang", "category": "indexability", "severity": "Low", "url": url})

        if is_self_canonical:
            if not soup.find('meta', attrs={'name': 'viewport'}):
                issues.append({"id": "missing_viewport", "category": "technical", "severity": "Critical", "url": url})
            
            if not soup.find('script', type='application/ld+json'):
                 path = urlparse(url).path.lower()
                 rec = "BreadcrumbList"
                 if path in ["/", ""]: rec = "Organization/WebSite"
                 elif any(x in path for x in ["product", "shop"]): rec = "Product"
                 elif any(x in path for x in ["blog", "news"]): rec = "Article"
                 issues.append({"id": "missing_jsonld", "category": "technical", "severity": "Medium", "url": url, "args": [rec]})

            if '_' in url: issues.append({"id": "url_underscore", "category": "technical", "severity": "Low", "url": url})
            if any(c.isupper() for c in urlparse(url).path): issues.append({"id": "url_uppercase", "category": "technical", "severity": "Medium", "url": url})
            
            if soup.find('a', href=lambda x: x and x.lower().startswith('javascript:')):
                issues.append({"id": "js_links", "category": "access", "severity": "High", "url": url}) 

            imgs = soup.find_all('img')
            missing_alt = 0
            bad_alt = 0
            cls_risk = 0
            for img in imgs:
                alt = img.get('alt', '').strip()
                if not alt: missing_alt += 1
                elif len(alt) < 3 or any(x in alt.lower() for x in ["image", "photo", "img"]): bad_alt += 1
                if not img.get('width') or not img.get('height'): cls_risk += 1
            
            if missing_alt > 0: issues.append({"id": "missing_alt", "category": "image_ux", "severity": "Medium", "url": url})
            if bad_alt > 0: issues.append({"id": "alt_bad_quality", "category": "image_ux", "severity": "Low", "url": url})
            if cls_risk > 0: issues.append({"id": "cls_risk", "category": "cwv_performance", "severity": "Medium", "url": url})

            links = soup.find_all('a', href=True)
            bad_anchors = ["click here", "read more", "more"]
            if any(a.get_text().strip().lower() in bad_anchors for a in links):
                issues.append({"id": "anchor_bad_quality", "category": "access", "severity": "Low", "url": url})
            
            if not title: 
                issues.append({"id": "missing_title", "category": "content", "severity": "High", "url": url})
            else:
                px_w = estimate_pixel_width(title)
                if px_w < 200:
                    issues.append({"id": "short_title", "category": "content", "severity": "Medium", "url": url, "evidence": title, "args": [int(px_w)]})
                elif px_w > 600:
                    issues.append({"id": "long_title", "category": "content", "severity": "Low", "url": url, "evidence": title, "args": [int(px_w)]})

            if not desc_content: 
                issues.append({"id": "missing_desc", "category": "content", "severity": "High", "url": url})
            else:
                px_w_d = estimate_pixel_width(desc_content)
                if px_w_d < 400:
                    issues.append({"id": "short_desc", "category": "content", "severity": "Low", "url": url, "evidence": desc_content, "args": [int(px_w_d)]})

            if not h1_content: issues.append({"id": "missing_h1", "category": "content", "severity": "High", "url": url})

            if (title and "not found" in title.lower()) or (soup.find('h1') and "not found" in soup.find('h1').get_text().lower()):
                issues.append({"id": "soft_404", "category": "access", "severity": "Critical", "url": url})
        
        if baidu_mode:
            keywords = soup.find('meta', attrs={'name': 'keywords'})
            if not keywords or not keywords.get('content', '').strip():
                 issues.append({"id": "missing_keywords", "category": "content", "severity": "Medium", "url": url})
            
            if "hm.baidu.com" not in str(soup):
                 issues.append({"id": "missing_baidu_stats", "category": "technical", "severity": "Low", "url": url})
            
            if not soup.find('meta', attrs={'name': 'applicable-device'}):
                 issues.append({"id": "missing_applicable_device", "category": "technical", "severity": "Medium", "url": url})
            
            has_no_transform = False
            for meta in soup.find_all('meta'):
                if meta.get('http-equiv', '').lower() == 'cache-control' and 'no-transform' in meta.get('content', '').lower():
                    has_no_transform = True
                    break
            if not has_no_transform:
                 issues.append({"id": "missing_no_transform", "category": "technical", "severity": "Medium", "url": url})
            
            page_text = soup.get_text()
            if "ICP备" not in page_text and "ICP证" not in page_text:
                 issues.append({"id": "missing_icp", "category": "technical", "severity": "High", "url": url})
            
            chinese_chars = len(re.findall(r'[\u4e00-\u9fa5]', page_text))
            total_chars = len(page_text.strip())
            if total_chars > 200 and (chinese_chars / total_chars) < 0.05:
                 issues.append({"id": "content_not_chinese", "category": "content", "severity": "Medium", "url": url})


    return {
        "URL": url, 
        "Status": status, 
        "Title": title, 
        "Description": desc_content,
        "H1": h1_content,
        "Canonical": can_url,
        "Content_Hash": hashlib.md5(soup.get_text().encode('utf-8')).hexdigest()
    }, issues

def crawl_website(start_url, max_pages, lang, manual_robots, manual_sitemaps, psi_key, list_url=None, detail_url=None, check_robots=True, crawl_sitemap=True, allow_sub=False, allow_outside=False, manual_pages=None, baidu_mode=False):
    visited = set()
    seen_hashes = {} 
    seen_urls = set()
    
    queue = [start_url]
    seen_urls.add(start_url)
    if list_url and is_valid_url(list_url):
         queue.append(list_url)
         seen_urls.add(list_url)
    if detail_url and is_valid_url(detail_url):
         queue.append(detail_url)
         seen_urls.add(detail_url)

    if manual_pages:
        for p in manual_pages:
            if is_valid_url(p) and p not in seen_urls:
                queue.append(p)
                seen_urls.add(p)

    results_data = []
    all_issues = []
    first_error = None
    target_domain = None
    
    start_netloc = urlparse(start_url).netloc.replace('www.', '')
    start_path = urlparse(start_url).path
    if not start_path.endswith('/'): start_path += '/'
    
    def clean_url(u): return u.split('?')[0].split('#')[0]

    progress_bar = st.progress(0, text="Initializing...")
    sitemap_has_hreflang = False
    
    try:
        site_issues, sitemap_has_hreflang = check_site_level_assets(
            start_url, lang, check_robots, crawl_sitemap, manual_sitemaps, baidu_mode
        )
        all_issues.extend(site_issues)
        st.session_state['sitemap_hreflang_found'] = sitemap_has_hreflang
    except Exception as e:
        pass

    if psi_key:
        with st.spinner(TRANSLATIONS[lang]["psi_fetching"].format("Pages")):
            targets = [("Home", start_url)]
            if list_url and is_valid_url(list_url): targets.append(("List", list_url))
            if detail_url and is_valid_url(detail_url): targets.append(("Detail", detail_url))
            
            for label, t_url in targets:
                cwv_data = fetch_psi_data(t_url, psi_key)
                if cwv_data and "error" not in cwv_data:
                    if label == "Home": st.session_state['cwv_data'] = cwv_data
                    all_issues.extend(check_cwv_issues(cwv_data, t_url, label=f"({label})"))

    count = 0
    headers = get_browser_headers()
    
    while queue and count < max_pages:
        url = queue.pop(0)
        visited.add(url)
        
        if any(x in url.lower() for x in ['/login', '/signin', '/admin', '/cart', '/account']):
            continue

        count += 1
        progress_bar.progress(int(count/max_pages*100), text=f"Crawling ({count}/{max_pages}): {url}")
        time.sleep(0.1)
        
        try:
            response = requests.get(url, headers=headers, timeout=10, allow_redirects=True, verify=False)
            current_url = response.url 
            
            if count == 1 and url == start_url:
                 start_netloc = urlparse(current_url).netloc.replace('www.', '')

            final_status = response.status_code

            if response.history:
                chain_list = [r.url for r in response.history] + [current_url]
                origin_netloc = urlparse(chain_list[0]).netloc.replace('www.', '')
                chain_display_parts = []
                for u in chain_list:
                    u_obj = urlparse(u)
                    u_netloc = u_obj.netloc.replace('www.', '')
                    
                    if u_netloc != origin_netloc:
                        chain_display_parts.append(u) # Full URL for cross-domain
                    else:
                        p = u_obj.path
                        if not p: p = "/"
                        chain_display_parts.append(p) # Path for same domain

                chain_str = " -> ".join(chain_display_parts)
                all_issues.append({"id": "http_3xx", "category": "access", "severity": "Medium", "url": url, "args": [chain_str]})

            if final_status >= 400:
                is_5xx = final_status >= 500
                all_issues.append({"id": "http_5xx" if is_5xx else "http_4xx", "category": "access", "severity": "Critical" if is_5xx else "High", "url": url, "args": [str(final_status)]})

            content_type = response.headers.get('Content-Type', '').lower()
            if 'text/html' in content_type:
                # Double check for login via content
                if 'type="password"' in response.text.lower():
                     continue # Skip login page content check

                page_data, page_issues = analyze_page(current_url, response.content, final_status, sitemap_has_hreflang, baidu_mode)
                
                # Deduplication & Data Storage
                if final_status == 200:
                    current_hash = page_data['Content_Hash']
                    current_canonical = page_data['Canonical']
                    current_clean = clean_url(current_url)
                    
                    if current_hash in seen_hashes:
                        original_url = seen_hashes[current_hash]
                        # Fix: Check if URL is actually different (avoid self-duplicate flagging)
                        if current_url != original_url and not (current_canonical and current_canonical != current_url):
                            all_issues.append({
                                "id": "duplicate", "category": "indexability", 
                                "severity": "High", "url": current_url, 
                                "meta": original_url # Raw URL
                            })
                    else:
                        seen_hashes[current_hash] = current_url

                results_data.append(page_data)
                all_issues.extend(page_issues)
                
                soup = BeautifulSoup(response.content, 'html.parser')
                for a in soup.find_all('a', href=True):
                    # Filter: No Fragment
                    raw_link = urljoin(current_url, a['href'])
                    link = raw_link.split('#')[0] 
                    
                    # Enhanced Filtering Logic
                    link_parsed = urlparse(link)
                    link_netloc = link_parsed.netloc.replace('www.', '')
                    link_path = link_parsed.path

                    # Check Domain
                    is_internal = False
                    if not link_netloc: is_internal = True # Relative
                    elif allow_sub:
                        is_internal = link_netloc.endswith(start_netloc) # Any subdomain
                    else:
                        is_internal = link_netloc == start_netloc # Strict match

                    # Check Path
                    path_ok = True
                    if not allow_outside:
                        if not link_path.startswith(start_path): path_ok = False
                    
                    if is_internal and path_ok and link not in seen_urls:
                        if not any(link.lower().endswith(ext) for ext in ['.jpg', '.png', '.pdf', '.zip', '.css', '.js', '.json', '.xml']):
                            seen_urls.add(link)
                            queue.append(link)
            else:
                if count == 1: first_error = f"Content type: {content_type}"
        except Exception as e:
            if count == 1: first_error = str(e)
            pass
    
    progress_bar.empty()
    if not results_data and first_error: return None, None, first_error
    return results_data, all_issues, None

# --- Level 7: 全局 PPT 绘图函数 ---
def set_font(font_obj, size, bold=False, color=None, lang="zh"):
    font_obj.size = Pt(size)
    font_obj.name = 'Microsoft YaHei' if lang == "zh" else 'Arial'
    font_obj.bold = bold
    if color: font_obj.color.rgb = color

def draw_cwv_gauge(slide, metric, value, thresholds):
    good, poor = thresholds
    total_w = 4
    x_start = 7.2
    y_pos = 4.8
    h = 0.4

    r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_start), Inches(y_pos), Inches(total_w/3), Inches(h))
    r1.fill.solid()
    r1.fill.fore_color.rgb = RGBColor(12, 206, 107)
    r1.line.fill.background()

    r2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_start + total_w/3), Inches(y_pos), Inches(total_w/3), Inches(h))
    r2.fill.solid()
    r2.fill.fore_color.rgb = RGBColor(255, 164, 0)
    r2.line.fill.background()
    
    r3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_start + 2*total_w/3), Inches(y_pos), Inches(total_w/3), Inches(h))
    r3.fill.solid()
    r3.fill.fore_color.rgb = RGBColor(255, 78, 66)
    r3.line.fill.background()
    
    def add_label(x_offset, text):
        tb = slide.shapes.add_textbox(Inches(x_start + x_offset - 0.2), Inches(y_pos + 0.40), Inches(0.5), Inches(0.3))
        tb.margin_top = 0
        p = tb.text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(80, 80, 80)

    add_label(0, "0")
    add_label(total_w/3, str(good))
    add_label(2*total_w/3, str(poor))

    pos = 0
    if value <= good:
        pos = (value / good) * (total_w/3)
    elif value <= poor:
        pos = (total_w/3) + ((value - good) / (poor - good)) * (total_w/3)
    else:
        cap = poor * 1.5
        normalized = min(value, cap)
        pos = (2*total_w/3) + ((normalized - poor) / (cap - poor)) * (total_w/3)

    marker_x = x_start + pos
    
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(marker_x - 0.1), Inches(y_pos - 0.2), Inches(0.2), Inches(0.2))
    tri.rotation = 180
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(50, 50, 50)
    tri.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(marker_x - 0.5), Inches(y_pos - 0.8), Inches(1), Inches(0.3))
    p = tb.text_frame.add_paragraph()
    p.text = f"{value}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def draw_serp_preview(slide, issue_id, issue_title, evidence, url, txt, lang="zh"):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4), Inches(5.8), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(220, 220, 220)
    
    label = slide.shapes.add_textbox(Inches(7), Inches(3.6), Inches(4), Inches(0.3))
    p = label.text_frame.add_paragraph()
    
    if "favicon" in issue_id:
        p.text = txt["visual_sim_title"]
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(4.3), Inches(0.25), Inches(0.25))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(200, 200, 200) 
        l1 = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(7.2), Inches(4.3), Inches(0.25), Inches(0.25))
        l1.line.color.rgb = RGBColor(150, 150, 150)
        l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.325), Inches(4.3), Inches(7.325), Inches(4.55))
        l2.line.color.rgb = RGBColor(150, 150, 150)

        tb = slide.shapes.add_textbox(Inches(7.5), Inches(4.25), Inches(4), Inches(0.4))
        p2 = tb.text_frame.add_paragraph()
        p2.text = urlparse(url).netloc
        set_font(p2.font, 12, False, RGBColor(32, 33, 36), lang)
        
        tb_t = slide.shapes.add_textbox(Inches(7.2), Inches(4.6), Inches(5), Inches(0.4))
        p_t = tb_t.text_frame.add_paragraph()
        p_t.text = "Page Title Example"
        set_font(p_t.font, 16, False, RGBColor(26, 13, 171), lang)
        
        tb_d = slide.shapes.add_textbox(Inches(7.2), Inches(5.0), Inches(5), Inches(0.4))
        p_d = tb_d.text_frame.add_paragraph()
        p_d.text = "This simulates a missing favicon result on Google mobile SERP."
        set_font(p_d.font, 12, False, RGBColor(80, 80, 80), lang)

    elif "alt" in issue_id:
        p.text = "Screen Reader View:"
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        
        img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(4.2), Inches(1.5), Inches(1.0))
        img_box.line.color.rgb = RGBColor(100, 100, 100)
        img_box.fill.background() 
        
        bubble_text = '<img src="..." />'
        color = RGBColor(200, 0, 0)
        if "quality" in issue_id:
            bubble_text = '<img src="..." alt="image" />'
            color = RGBColor(255, 165, 0)
        
        callout = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2, Inches(8.8), Inches(4.2), Inches(3.5), Inches(0.8))
        callout.fill.solid()
        callout.fill.fore_color.rgb = RGBColor(240, 240, 240)
        callout.text_frame.text = bubble_text
        callout.text_frame.paragraphs[0].font.color.rgb = color
            
    elif "lcp" in issue_id:
        p.text = txt["cwv_sim_title"] + " LCP"
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        val = float(re.findall(r"(\d+\.\d+)", evidence)[0]) if re.findall(r"(\d+\.\d+)", evidence) else 3.0
        draw_cwv_gauge(slide, "LCP", val, (2.5, 4.0))

    elif "inp" in issue_id:
        p.text = txt["cwv_sim_title"] + " INP"
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        val = float(re.findall(r"(\d+)", evidence)[0]) if re.findall(r"(\d+)", evidence) else 300
        draw_cwv_gauge(slide, "INP", val, (200, 500))

    elif "cls" in issue_id:
        p.text = txt["cwv_sim_title"] + " CLS"
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        val = float(re.findall(r"(\d+\.\d+)", evidence)[0]) if re.findall(r"(\d+\.\d+)", evidence) else 0.2
        draw_cwv_gauge(slide, "CLS", val, (0.1, 0.25))

    elif "fcp" in issue_id:
        p.text = txt["cwv_sim_title"] + " FCP"
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        val = float(re.findall(r"(\d+\.\d+)", evidence)[0]) if re.findall(r"(\d+\.\d+)", evidence) else 2.0
        thresholds = (1.0, 2.0) if "baidu" in issue_id else (1.8, 3.0)
        draw_cwv_gauge(slide, "FCP", val, thresholds)
             
    elif "3xx" in issue_id:
        p.text = "Redirect Flow:"
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        
        parts = evidence.split(' -> ')
        if len(parts) > 1:
            x = 7.2
            display_parts = parts[:3]
            if len(parts) > 3: display_parts[-1] = "Final"

            for i, part in enumerate(display_parts):
                box_w = 2.0 if len(part) > 20 else 1.5
                b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.5), Inches(box_w), Inches(0.6))
                b.fill.solid()
                b.fill.fore_color.rgb = RGBColor(240, 240, 240) if i < len(display_parts)-1 else RGBColor(220, 255, 220)
                b.text_frame.text = part
                b.text_frame.paragraphs[0].font.size = Pt(9)
                b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                
                if i < len(display_parts) - 1:
                    ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+box_w), Inches(4.7), Inches(0.3), Inches(0.2))
                    ar.fill.solid()
                    ar.fill.fore_color.rgb = RGBColor(100, 100, 100)
                    x += (box_w + 0.3)
            
    else:
        p.text = txt["serp_sim_title"]
        set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)
        
        tf = box.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p_serp = tf.add_paragraph()
        p_serp.text = f"{urlparse(url).netloc} › ..."
        set_font(p_serp.font, 12, False, RGBColor(32, 33, 36), lang)
        
        if "short_desc" in issue_id or "missing_desc" in issue_id:
            p_serp = tf.add_paragraph()
            p_serp.space_before = Pt(5)
            p_serp.text = evidence[:60] + "..." if evidence else "Title of the page"
            set_font(p_serp.font, 18, False, RGBColor(26, 13, 171), lang)
            
            p_serp = tf.add_paragraph()
            p_serp.space_before = Pt(3)
            if "missing" in issue_id:
                p_serp.text = "No description available in code..."
            else:
                p_serp.text = evidence 
            set_font(p_serp.font, 14, False, RGBColor(77, 81, 86), lang)
        
        elif "long_title" in issue_id:
            p_serp = tf.add_paragraph()
            p_serp.space_before = Pt(5)
            p_serp.text = evidence[:55] + " ..."
            set_font(p_serp.font, 18, False, RGBColor(26, 13, 171), lang)

            p_serp = tf.add_paragraph()
            p_serp.space_before = Pt(3)
            p_serp.text = "The meta description of the page would appear here..."
            set_font(p_serp.font, 14, False, RGBColor(77, 81, 86), lang)

        else:
            p_serp = tf.add_paragraph()
            p_serp.space_before = Pt(5)
            p_serp.text = evidence if evidence else "Untitled Page"
            set_font(p_serp.font, 18, False, RGBColor(26, 13, 171), lang)

def draw_code_preview(slide, txt, lang="zh"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(4), Inches(5.8), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box.line.color.rgb = RGBColor(200, 200, 200)
    tf = box.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.add_paragraph()
    p.text = '<a href="javascript:void(0)">\n  Click Here\n</a>'
    set_font(p.font, 14, True, RGBColor(200, 0, 0), lang)
    
    label = slide.shapes.add_textbox(Inches(7), Inches(3.6), Inches(4), Inches(0.3))
    p = label.text_frame.add_paragraph()
    p.text = txt["code_sim_title"]
    set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)

def draw_hreflang_preview(slide, url, missing_type, txt, lang="zh"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(4), Inches(5.8), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box.line.color.rgb = RGBColor(200, 200, 200)
    
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    p = tf.add_paragraph()
    p.text = "<!-- Correct Implementation -->"
    set_font(p.font, 10, False, RGBColor(128, 128, 128), lang)
    
    p = tf.add_paragraph()
    clean_url = url.split('?')[0][:40] + "..."
    if "default" in missing_type:
        p.text = f'<link rel="alternate" hreflang="x-default" href="{clean_url}" />'
        set_font(p.font, 11, True, RGBColor(200, 0, 0), lang)
    else:
        p.text = f'<link rel="alternate" hreflang="en" href="{clean_url}" />\n<link rel="alternate" hreflang="fr" href="..." />'
        set_font(p.font, 11, False, RGBColor(0, 0, 128), lang)

    label = slide.shapes.add_textbox(Inches(7), Inches(3.6), Inches(4), Inches(0.3))
    p = label.text_frame.add_paragraph()
    p.text = txt["code_sim_title"]
    set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)

def draw_rich_snippet_preview(slide, url, txt, lang="zh"):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4), Inches(5.8), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(220, 220, 220)
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.add_paragraph()
    p.text = f"{urlparse(url).netloc} › product"
    set_font(p.font, 12, False, RGBColor(32, 33, 36), lang)
    
    p = tf.add_paragraph()
    p.space_before = Pt(5)
    p.text = "Product Name Example - Best Choice"
    set_font(p.font, 18, False, RGBColor(26, 13, 171), lang)
    
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    p.text = "★★★★★ Rating: 4.8 · $199.00 · In stock"
    set_font(p.font, 12, False, RGBColor(231, 113, 27), lang)
    
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    p.text = "This is a rich result enabled by Schema..."
    set_font(p.font, 14, False, RGBColor(77, 81, 86), lang)
    
    label = slide.shapes.add_textbox(Inches(7), Inches(3.6), Inches(4), Inches(0.3))
    p = label.text_frame.add_paragraph()
    p.text = txt["rich_sim_title"]
    set_font(p.font, 12, True, RGBColor(100, 100, 100), lang)

def create_styled_pptx(slides_data, lang):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    txt = TRANSLATIONS[lang] 
    
    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(18, 52, 86)
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    p = title.text_frame.add_paragraph()
    p.text = txt["ppt_cover_title"]
    p.alignment = PP_ALIGN.CENTER
    set_font(p.font, 54, True, RGBColor(255, 255, 255), lang)
    
    sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    p = sub.text_frame.add_paragraph()
    p.text = txt["ppt_cover_sub"]
    p.alignment = PP_ALIGN.CENTER
    set_font(p.font, 24, False, RGBColor(200, 200, 200), lang)

    # Slides
    for s in slides_data:
        t_data = get_translated_text(s['id'], lang, s.get('args'))
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header
        h_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
        h_shape.fill.solid()
        h_shape.fill.fore_color.rgb = RGBColor(240, 242, 246)
        h_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        p = h_box.text_frame.add_paragraph()
        p.text = t_data['title']
        set_font(p.font, 32, True, RGBColor(50, 50, 50), lang)
        
        sev_color = RGBColor(220, 53, 69) if s['severity'] == "Critical" else RGBColor(253, 126, 20)
        sev_box = slide.shapes.add_textbox(Inches(11), Inches(0.35), Inches(2), Inches(0.5))
        p = sev_box.text_frame.add_paragraph()
        p.text = s['severity']
        p.alignment = PP_ALIGN.CENTER
        set_font(p.font, 18, True, sev_color, lang)
        
        # Category
        cat_key = f"cat_{s['category']}"
        cat_label = txt.get(cat_key, s['category'].upper())
        cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.4))
        p = cat_box.text_frame.add_paragraph()
        p.text = cat_label
        set_font(p.font, 14, True, RGBColor(0, 102, 204), lang)

        # Desc
        d_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(0.5))
        p = d_title.text_frame.add_paragraph()
        p.text = txt["ppt_desc"]
        set_font(p.font, 18, True, RGBColor(30, 30, 30), lang)
        
        d_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(6), Inches(1.2))
        tf = d_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.add_paragraph()
        p.text = t_data['desc']
        set_font(p.font, 14, False, RGBColor(80, 80, 80), lang)
        
        # Impact
        i_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(6), Inches(0.5))
        p = i_title.text_frame.add_paragraph()
        p.text = txt["ppt_business_impact"]
        set_font(p.font, 18, True, RGBColor(30, 30, 30), lang)
        
        i_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(6), Inches(1.2))
        tf = i_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.add_paragraph()
        p.text = t_data['impact']
        set_font(p.font, 14, False, RGBColor(220, 53, 69), lang)

        # Suggestion
        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.5))
        s_bg.fill.solid()
        s_bg.fill.fore_color.rgb = RGBColor(230, 244, 234)
        s_bg.line.color.rgb = RGBColor(40, 167, 69)
        s_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(11.9), Inches(1.3))
        tf = s_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.add_paragraph()
        p.text = txt["ppt_slide_sugg_title"]
        set_font(p.font, 16, True, RGBColor(21, 87, 36), lang)
        p = tf.add_paragraph()
        p.text = t_data['suggestion']
        set_font(p.font, 14, False, RGBColor(21, 87, 36), lang)

        # Right Column
        ex_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(0.5))
        p = ex_title.text_frame.add_paragraph()
        p.text = txt["ppt_slide_ex_title"]
        set_font(p.font, 18, True, RGBColor(30, 30, 30), lang)
        
        ex_box = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5.8), Inches(1.5)) 
        tf = ex_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        for idx, url in enumerate(s['examples'][:4]): 
            p = tf.add_paragraph()
            if s['id'] == 'duplicate' and "Duplicate Group:" in url:
                 parts = url.split("\n")
                 p.text = f"Group {idx+1}:\n   • {parts[1].replace('- ', '').strip()}\n   • {parts[2].replace('- ', '').strip()}"
                 set_font(p.font, 10, False, RGBColor(80, 80, 80), lang)
            else:
                 p.text = f"• {url}"
                 set_font(p.font, 11, False, RGBColor(0, 102, 204), lang)
            p.space_after = Pt(6)

        # Visualization Logic
        is_serp = any(k in s['id'] for k in ["title", "desc", "favicon", "alt", "lcp", "inp", "cls", "3xx", "fcp"])
        is_rich = "jsonld" in s['id']
        is_code = "js_links" in s['id'] or "anchor" in s['id']
        is_hreflang = "hreflang" in s['id']
        is_cwv = any(k in s['id'] for k in ["lcp", "inp", "cls", "fcp", "risk"])
        is_img = "alt" in s['id'] or "favicon" in s['id']
        is_3xx = "3xx" in s['id'] 
        
        ev = s.get('example_evidence', '')
        ex_url = s['examples'][0] if s['examples'] else "example.com"
        if "Duplicate" in ex_url: ex_url = ex_url.split("\n")[1].replace("- ", "").strip()
        if "3xx" in s['id'] and s.get('args'): ev = s['args'][0]

        if is_code:
            draw_code_preview(slide, txt, lang)
        elif is_hreflang:
            type_str = s['id']
            if "invalid" in type_str and s.get('args'):
                type_str = f"invalid: {s['args'][0]}"
            draw_hreflang_preview(slide, ex_url, type_str, txt, lang)
        elif is_rich:
            draw_rich_snippet_preview(slide, ex_url, txt, lang)
        elif is_serp:
            draw_serp_preview(slide, s['id'], t_data['title'], ev, ex_url, txt, lang)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# --- 7. UI Logic (No indentation) ---
if 'audit_data' not in st.session_state: st.session_state['audit_data'] = None
if 'audit_issues' not in st.session_state: st.session_state['audit_issues'] = []
if 'language' not in st.session_state: st.session_state['language'] = "zh"
if 'cwv_data' not in st.session_state: st.session_state['cwv_data'] = None
if 'sitemap_hreflang_found' not in st.session_state: st.session_state['sitemap_hreflang_found'] = False

lang = st.session_state['language']
ui = TRANSLATIONS[lang]

with st.sidebar:
    st.title(ui["sidebar_title"])
    st.caption(ui["sidebar_caption"])
    st.divider()
    
    sl = st.radio(ui["lang_label"], ["中文", "English"], index=0 if lang=="zh" else 1)
    if (sl == "中文" and lang == "en") or (sl == "English" and lang == "zh"):
        st.session_state['language'] = "zh" if sl == "中文" else "en"
        st.rerun()

    st.divider()
    opts = ui["nav_options"]
    keys = ["input", "dashboard", "matrix", "ppt"]
    sel = st.radio(ui["nav_label"], opts)
    menu_key = keys[opts.index(sel)]
    
    st.divider()
    if st.session_state['audit_data']:
        st.success(ui["cache_info"].format(len(st.session_state['audit_data'])))
        st.markdown(f"**{ui['sitemap_status_title']}**")
        if st.session_state['sitemap_hreflang_found']: st.caption(ui["sitemap_found_href"])
        else: st.caption(ui["sitemap_no_href"])
        
        if st.button(ui["clear_data"]):
            st.session_state['audit_data'] = None
            st.session_state['audit_issues'] = []
            st.session_state['cwv_data'] = None
            st.rerun()

if menu_key == "input":
    st.header(ui["input_header"])
    st.info(ui["input_info"])
    c1, c2 = st.columns([3, 1])
    with c1: target_url = st.text_input(ui["input_label"], placeholder=ui["input_placeholder"])
    with c2: max_pages = st.number_input(ui.get("max_pages_label", "Max Pages"), min_value=1, max_value=1000, value=100)
    
    with st.expander(ui.get("adv_settings", "Advanced")):
        allow_sub = st.checkbox(ui["allow_subdomains_label"], value=False)
        allow_out = st.checkbox(ui["allow_outside_folder_label"], value=False)
        check_robots_flag = st.checkbox(ui["check_robots_label"], value=True)
        crawl_sitemap_flag = st.checkbox(ui["crawl_sitemap_label"], value=True)
        baidu_mode_flag = st.checkbox(ui["baidu_mode_label"], value=False)
        manual_sitemaps_text = st.text_area(ui.get("manual_sitemaps", "Manual Sitemaps"), placeholder="https://example.com/sitemap.xml")
        manual_sitemaps = [s.strip() for s in manual_sitemaps_text.split('\n') if s.strip()]
        manual_pages_text = st.text_area(ui.get("manual_pages_label", "Manual Pages"), placeholder="https://example.com/page1")
        manual_pages = [s.strip() for s in manual_pages_text.split('\n') if s.strip()]
        sitemap_content_text = st.text_area(ui.get("sitemap_content_label", "Paste Sitemap Content"), height=150)

    
    with st.expander(ui.get("psi_settings", "Google PSI")):
        psi_key = st.text_input(ui.get("psi_api_key_label", "API Key"), type="password", help=ui.get("psi_api_help", ""))
        psi_list_url = st.text_input(ui.get("psi_list_url_label", "List URL"))
        psi_detail_url = st.text_input(ui.get("psi_detail_url_label", "Detail URL"))
        st.caption(ui["psi_get_key"])

    if st.button(ui["start_btn"], type="primary"):
        if not target_url or not is_valid_url(target_url): 
            st.error(ui["error_url"])
        else:
            with st.spinner(ui["spinner_crawl"].format(max_pages)):
                # Handle pasted sitemap content
                if sitemap_content_text:
                    extracted_urls = re.findall(r'<loc>\s*(https?://[^<]+)\s*</loc>', sitemap_content_text)
                    if manual_pages is None: manual_pages = []
                    manual_pages.extend(extracted_urls)

                data, issues, error_msg = crawl_website(
                    target_url, max_pages, lang, None, manual_sitemaps, psi_key, 
                    psi_list_url, psi_detail_url, check_robots_flag, crawl_sitemap_flag,
                    allow_sub, allow_out, manual_pages, baidu_mode_flag
                )
                if not data:
                    st.error(ui["error_no_data"].format(error_msg or "Unknown Error"))
                else:
                    st.session_state['audit_data'] = data
                    st.session_state['audit_issues'] = issues
                    st.success(ui["success_audit"].format(len(data)))
                    st.balloons()

elif menu_key == "dashboard":
    st.header(ui["dashboard_header"])
    if not st.session_state['audit_data']: st.warning(ui["warn_no_data"])
    else:
        if st.session_state.get('cwv_data'):
            c = st.session_state['cwv_data']
            st.subheader(ui["cwv_title"])
            st.caption(ui["cwv_source"])
            c1, c2, c3, c4 = st.columns(4)
            def metric_color(val, good, poor):
                if val <= good: return "normal"
                if val >= poor: return "inverse"
                return "off"
            c1.metric("LCP (Loading)", f"{c['LCP']:.2f}s", delta_color=metric_color(c['LCP'], 2.5, 4.0))
            c2.metric("CLS (Visual)", f"{c['CLS']:.3f}", delta_color=metric_color(c['CLS'], 0.1, 0.25))
            c3.metric("INP (Interact)", f"{c['INP']}ms", delta_color=metric_color(c['INP'], 200, 500))
            c4.metric("FCP", f"{c['FCP']:.2f}s")
            st.divider()

        df = pd.DataFrame(st.session_state['audit_data'])
        issues = st.session_state['audit_issues']
        total = len(issues)
        score = max(0, 100 - int(total * 0.5))
        critical = len([i for i in issues if i['severity'] == 'Critical'])
        
        k1, k2, k3, k4 = st.columns(4)
        k1.metric(ui["kpi_health"], f"{score}/100")
        k2.metric(ui["kpi_pages"], str(len(df)))
        k3.metric(ui["kpi_issues"], str(total), delta_color="inverse")
        k4.metric(ui["kpi_critical"], str(critical), delta_color="inverse")
        
        st.divider()
        c1, c2 = st.columns(2)
        with c1:
            st.subheader(ui["chart_issues"])
            if issues:
                issue_counts = pd.DataFrame(issues)['id'].value_counts().reset_index()
                issue_counts.columns = ['id', 'count']
                issue_counts['name'] = issue_counts['id'].apply(lambda x: get_translated_text(x, lang)['title'])
                st.bar_chart(issue_counts.set_index('name'))
            else: st.info(ui["chart_no_issues"])
        with c2:
            st.subheader(ui["chart_status"])
            if not df.empty: st.bar_chart(df['Status'].value_counts())

elif menu_key == "matrix":
    st.header(ui["matrix_header"])
    if not st.session_state['audit_data']: st.warning(ui["warn_no_data"])
    else:
        df = pd.DataFrame(st.session_state['audit_data'])
        st.dataframe(df, use_container_width=True)
        st.download_button(ui["download_csv"], df.to_csv().encode('utf-8'), "audit.csv")

elif menu_key == "ppt":
    st.header(ui["ppt_header"])
    if not st.session_state['audit_issues']: st.warning(ui["warn_no_data"])
    else:
        raw = st.session_state['audit_issues']
        grouped = {}
        for i in raw:
            iid = i['id']
            if iid not in grouped:
                grouped[iid] = {
                    "id": iid, "category": i['category'], "severity": i['severity'],
                    "count": 0, "examples": [], "args": i.get('args', []),
                    "example_evidence": i.get("evidence", "")
                }
            grouped[iid]['count'] += 1
            if len(grouped[iid]['examples']) < 5:
                if iid == "duplicate" and "meta" in i:
                     # Clean grouping for duplicate
                     grouped[iid]['examples'].append(f"Duplicate Group:\n- {i['url']}\n- {i['meta']}")
                else:
                     grouped[iid]['examples'].append(i['url'])
        
        slides = sorted(list(grouped.values()), key=lambda x: (
            CATEGORY_ORDER.index(x['category']),
            get_issue_priority(x['id']),
            SEVERITY_ORDER.get(x['severity'], 3)
        ))
        
        st.write(f"### {ui['ppt_download_header']}")
        st.info(ui["ppt_info"])
        if st.button(ui["ppt_btn"]):
            with st.spinner("Generating..."):
                f = create_styled_pptx(slides, lang)
                st.download_button(ui["ppt_btn"], f, f"seo_audit_{lang}.pptx")
        
        if 'slide_index' not in st.session_state: st.session_state.slide_index = 0
        if st.session_state.slide_index >= len(slides): st.session_state.slide_index = 0
        
        s = slides[st.session_state.slide_index]
        t_data = get_translated_text(s['id'], lang, s['args'])
        
        with st.container(border=True):
            st.caption(f"📂 {ui.get('cat_'+s['category'], s['category'])}")
            st.markdown(f"### {ui['ppt_slide_title']} {t_data['title']}")
            
            c1, c2 = st.columns([1, 1])
            with c1:
                color = "red" if s['severity'] == "Critical" else "orange"
                st.markdown(f"**{ui['ppt_severity']}** :{color}[{s['severity']}]")
                st.markdown(f"**{ui['ppt_impact']}** {ui['ppt_impact_desc'].format(s['count'])}")
                
                st.markdown(f"**{ui['ppt_desc']}**")
                st.write(t_data['desc'])
                
                st.markdown(f"**{ui['ppt_business_impact']}**") 
                st.error(t_data['impact']) 
                
                st.info(f"{ui['ppt_sugg']} {t_data['suggestion']}")
            with c2:
                # Visualization Logic
                is_serp = any(k in s['id'] for k in ["title", "desc", "favicon", "alt", "lcp", "inp", "cls", "3xx", "fcp"])
                is_rich = "jsonld" in s['id']
                is_code = "js_links" in s['id'] or "anchor" in s['id']
                is_hreflang = "hreflang" in s['id']
                is_cwv = any(k in s['id'] for k in ["lcp", "inp", "cls", "fcp", "risk"])
                is_img = "alt" in s['id'] or "favicon" in s['id']
                is_3xx = "3xx" in s['id'] 
                
                ev = s.get('example_evidence', '')
                ex_url = s['examples'][0] if s['examples'] else "example.com"
                if "Duplicate" in ex_url: ex_url = ex_url.split("\n")[1].replace("- ", "").strip()
                if "3xx" in s['id'] and s.get('args'): ev = s['args'][0]

                if is_code:
                    st.code('<link rel="alternate" ... />' if "hreflang" in s['id'] else '<a href="javascript:...">', language="html")
                elif is_hreflang:
                    type_str = s['id']
                    if "invalid" in type_str and s.get('args'):
                        type_str = f"invalid: {s['args'][0]}"
                    st.code(f'<link rel="alternate" hreflang="{type_str}" href="..." />', language="html")
                elif is_rich:
                    st.markdown(f"**{ui.get('rich_sim_title', 'Rich Result Preview')}**")
                    st.markdown("""
                     <div style="font-family: Arial, sans-serif; border: 1px solid #dfe1e5; border-radius: 8px; padding: 15px; background: white; box-shadow: 0 1px 6px rgba(32,33,36,0.28);">
                        <div style="font-size: 14px; color: #202124;">example.com <span style="color: #5f6368">› product</span></div>
                        <div style="font-size: 20px; color: #1a0dab; margin-top: 5px;">Best Product - High Quality</div>
                        <div style="color: #e7711b; font-size: 14px;">★★★★★ <span style="color:#70757a">Rating: 4.8 · $199.00 · In stock</span></div>
                        <div style="font-size: 14px; color: #4d5156; margin-top: 3px;">This is a rich result enabled by Schema...</div>
                     </div>
                     """, unsafe_allow_html=True)
                elif is_serp:
                    st.markdown(f"**{ui.get('serp_sim_title', 'SERP Preview')}**")
                    display_title = ev if ev else "Untitled Page"
                    if len(display_title) > 60: display_title = display_title[:55] + " ..."
                    st.markdown(f"""
                    <div style="font-family: Arial, sans-serif; border: 1px solid #dfe1e5; border-radius: 8px; padding: 15px; background: white; box-shadow: 0 1px 6px rgba(32,33,36,0.28);">
                        <div style="font-size: 14px; color: #202124;">{urlparse(ex_url).netloc} <span style="color: #5f6368">› ...</span></div>
                        <div style="font-size: 20px; color: #1a0dab; margin-top: 5px;">{display_title}</div>
                        <div style="font-size: 14px; color: #4d5156; margin-top: 3px;">
                            Please provide a meta description...
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                elif is_cwv or is_img or is_3xx:
                    st.warning(f"See PPT for full visual simulation of {t_data['title']}")
                
                st.markdown(f"**{ui['ppt_examples']}**")
                for ex in s['examples']:
                     if "Duplicate Group:" in ex:
                         parts = ex.split("\n")
                         st.markdown(f"- **Group:**\n  - `{parts[1].replace('- ', '').strip()}`\n  - `{parts[2].replace('- ', '').strip()}`")
                     else:
                         st.markdown(f"- `{ex}`")

        c1, c2, c3 = st.columns([1, 2, 1])
        if c1.button(ui["ppt_prev"]): 
            st.session_state.slide_index = max(0, st.session_state.slide_index-1)
            st.rerun()
        c2.markdown(f"<div style='text-align: center'>Slide {st.session_state.slide_index+1} / {len(slides)}</div>", unsafe_allow_html=True)
        if c3.button(ui["ppt_next"]):
            st.session_state.slide_index = min(len(slides)-1, st.session_state.slide_index+1)
            st.rerun()
